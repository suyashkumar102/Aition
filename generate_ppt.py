"""
Aition — Prototype Deck Generator
Run: pip install python-pptx && python generate_ppt.py
Output: Aition_Prototype_Deck.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ───────────────────────────────────────────────────────────────────
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x0D, 0x0D, 0x0D)
PURPLE     = RGBColor(0x5A, 0x3B, 0xFF)
PURPLE_LT  = RGBColor(0xA4, 0x8B, 0xFF)
RED        = RGBColor(0xDC, 0x26, 0x26)
GREEN      = RGBColor(0x05, 0x96, 0x69)
GRAY       = RGBColor(0x6B, 0x72, 0x80)
LIGHT_GRAY = RGBColor(0xF3, 0xF4, 0xF6)
DARK_GRAY  = RGBColor(0x37, 0x41, 0x51)
ACCENT     = RGBColor(0x7C, 0x3A, 0xED)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank

# ── helpers ───────────────────────────────────────────────────────────────────

def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE
    return s

def box(s, x, y, w, h, fill=None, line=None, line_w=Pt(1)):
    from pptx.util import Emu
    shape = s.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def txt(s, text, x, y, w, h, size=18, bold=False, color=BLACK,
        align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def accent_bar(s, color=PURPLE):
    """Left accent bar on content slides."""
    box(s, Inches(0), Inches(0), Inches(0.06), H, fill=color)

def slide_number(s, n, total=10):
    txt(s, f"{n} / {total}", W - Inches(1.2), H - Inches(0.45),
        Inches(1.1), Inches(0.35), size=10, color=GRAY, align=PP_ALIGN.RIGHT)

def section_tag(s, label, color=PURPLE):
    b = box(s, Inches(0.55), Inches(0.28), Inches(1.8), Inches(0.32),
            fill=color)
    txt(s, label.upper(), Inches(0.55), Inches(0.28), Inches(1.8), Inches(0.32),
        size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

def heading(s, text, y=Inches(0.75), size=32, color=BLACK):
    txt(s, text, Inches(0.55), y, Inches(12.2), Inches(0.7),
        size=size, bold=True, color=color)

def divider(s, y):
    box(s, Inches(0.55), y, Inches(12.2), Pt(1.5), fill=LIGHT_GRAY)

def bullet_block(s, items, x, y, w, h, size=14, color=DARK_GRAY, gap=Inches(0.38)):
    cy = y
    for item in items:
        prefix = "•  " if not item.startswith("–") else "   "
        txt(s, prefix + item.lstrip("–").strip(), x, cy, w, gap,
            size=size, color=color)
        cy += gap
    return cy

def card(s, x, y, w, h, title, body_lines, title_color=PURPLE,
         bg=LIGHT_GRAY, border=None):
    box(s, x, y, w, h, fill=bg, line=border or LIGHT_GRAY)
    txt(s, title, x + Inches(0.18), y + Inches(0.12), w - Inches(0.3),
        Inches(0.35), size=13, bold=True, color=title_color)
    cy = y + Inches(0.5)
    for line in body_lines:
        txt(s, "• " + line, x + Inches(0.18), cy,
            w - Inches(0.3), Inches(0.32), size=11, color=DARK_GRAY)
        cy += Inches(0.3)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════════════════
s1 = slide()

# Purple top band
box(s1, Inches(0), Inches(0), W, Inches(0.08), fill=PURPLE)

# Logo pill
box(s1, Inches(0.55), Inches(1.1), Inches(1.1), Inches(1.1),
    fill=PURPLE, line=PURPLE)
txt(s1, "A", Inches(0.55), Inches(1.1), Inches(1.1), Inches(1.1),
    size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Title
txt(s1, "Aition", Inches(1.85), Inches(1.15), Inches(10), Inches(0.85),
    size=52, bold=True, color=BLACK)
txt(s1, "Causal AI Fairness Engine", Inches(1.85), Inches(2.05),
    Inches(10), Inches(0.55), size=24, color=PURPLE, bold=True)

# Tagline
box(s1, Inches(0.55), Inches(2.85), Inches(11.5), Pt(1.5), fill=LIGHT_GRAY)
txt(s1,
    '"Your hiring model just passed every standard fairness test.\n'
    ' It rejected 81 qualified candidates anyway."',
    Inches(0.55), Inches(3.05), Inches(11.5), Inches(1.0),
    size=17, italic=True, color=DARK_GRAY)

# Three stat pills
for i, (val, lbl, col) in enumerate([
    ("3", "Proxy Paths Detected", PURPLE),
    ("81", "Candidates Affected", RED),
    ("0", "Found by Standard Tools", GREEN),
]):
    bx = Inches(0.55) + i * Inches(4.1)
    box(s1, bx, Inches(4.4), Inches(3.8), Inches(1.5),
        fill=LIGHT_GRAY, line=col)
    txt(s1, val, bx, Inches(4.5), Inches(3.8), Inches(0.75),
        size=40, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(s1, lbl, bx, Inches(5.2), Inches(3.8), Inches(0.45),
        size=12, color=GRAY, align=PP_ALIGN.CENTER)

# Footer
box(s1, Inches(0), H - Inches(0.55), W, Inches(0.55), fill=LIGHT_GRAY)
txt(s1, "Google Solution Challenge 2026 India  |  Unbiased AI Decision Track",
    Inches(0.55), H - Inches(0.48), Inches(9), Inches(0.38),
    size=11, color=GRAY)
txt(s1, "Prototype Deck", W - Inches(2.2), H - Inches(0.48),
    Inches(2.0), Inches(0.38), size=11, color=GRAY, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — BRIEF ABOUT THE SOLUTION
# ══════════════════════════════════════════════════════════════════════════════
s2 = slide()
accent_bar(s2)
section_tag(s2, "Slide 1")
slide_number(s2, 1)
heading(s2, "Brief About the Solution")
divider(s2, Inches(1.55))

txt(s2,
    "Aition is a Causal AI Fairness Engine that detects proxy discrimination "
    "— the hidden bias that passes every standard fairness test.",
    Inches(0.55), Inches(1.7), Inches(12.2), Inches(0.6),
    size=15, color=DARK_GRAY)

# Problem box
box(s2, Inches(0.55), Inches(2.45), Inches(5.8), Inches(3.8),
    fill=RGBColor(0xFF,0xF1,0xF1), line=RED)
txt(s2, "THE PROBLEM", Inches(0.75), Inches(2.6), Inches(5.4), Inches(0.35),
    size=11, bold=True, color=RED)
prob_items = [
    "Standard tools only measure output rate equality",
    "They CANNOT detect proxy discrimination",
    "A model can pass demographic parity and still",
    "  discriminate through neutral-looking variables",
    "10M+ automated decisions/day in India affected",
    "No tool surfaces the fairness impossibility theorem",
]
bullet_block(s2, prob_items, Inches(0.75), Inches(3.05),
             Inches(5.3), Inches(3.8), size=12, color=DARK_GRAY,
             gap=Inches(0.36))

# Arrow
txt(s2, "→", Inches(6.5), Inches(3.9), Inches(0.5), Inches(0.5),
    size=28, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

# Solution box
box(s2, Inches(7.1), Inches(2.45), Inches(5.65), Inches(3.8),
    fill=RGBColor(0xF5,0xF3,0xFF), line=PURPLE)
txt(s2, "AITION'S SOLUTION", Inches(7.3), Inches(2.6), Inches(5.2),
    Inches(0.35), size=11, bold=True, color=PURPLE)
sol_items = [
    "Builds a causal graph of how the model decides",
    "Traces proxy paths using do-calculus (DoWhy)",
    "Detects: age → graduation gap → hired",
    "Detects: SES → neighbourhood quality → hired",
    "Surfaces the fairness impossibility tradeoff",
    "Generates plain-language reports via Gemini AI",
]
bullet_block(s2, sol_items, Inches(7.3), Inches(3.05),
             Inches(5.2), Inches(3.8), size=12, color=DARK_GRAY,
             gap=Inches(0.36))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — OPPORTUNITIES & USP
# ══════════════════════════════════════════════════════════════════════════════
s3 = slide()
accent_bar(s3)
section_tag(s3, "Slide 2", color=ACCENT)
slide_number(s3, 2)
heading(s3, "Opportunities & USP")
divider(s3, Inches(1.55))

# Three columns
cols = [
    ("How is it different?", PURPLE, [
        "Only tool using causal graphs",
        "not just statistical correlation",
        "Detects proxy paths standard",
        "tools are blind to",
        "Two protected attributes",
        "audited simultaneously",
        "Interactive impossibility",
        "surface — not a single verdict",
    ]),
    ("How does it solve the problem?", RED, [
        "Builds DoWhy causal model",
        "from the dataset",
        "Applies backdoor criterion",
        "to find proxy paths",
        "Surgical debiasing removes",
        "only the biased components",
        "Gemini report explains it",
        "to non-technical stakeholders",
    ]),
    ("USP", GREEN, [
        "First tool to make the",
        "impossibility theorem",
        "interactive & actionable",
        "Plain-language audit report",
        "for HR directors, not just",
        "ML engineers",
        "Indian PDPB 2023 compliance",
        "documentation built-in",
    ]),
]

for i, (title, color, items) in enumerate(cols):
    bx = Inches(0.55) + i * Inches(4.2)
    box(s3, bx, Inches(1.75), Inches(3.95), Inches(4.8),
        fill=LIGHT_GRAY, line=color)
    txt(s3, title, bx + Inches(0.15), Inches(1.88),
        Inches(3.65), Inches(0.4), size=13, bold=True, color=color)
    box(s3, bx, Inches(2.3), Inches(3.95), Pt(1.5), fill=color)
    bullet_block(s3, items, bx + Inches(0.15), Inches(2.42),
                 Inches(3.65), Inches(4.8), size=11.5,
                 color=DARK_GRAY, gap=Inches(0.33))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — FEATURES
# ══════════════════════════════════════════════════════════════════════════════
s4 = slide()
accent_bar(s4)
section_tag(s4, "Slide 3", color=GREEN)
slide_number(s4, 3)
heading(s4, "Features Offered by the Solution")
divider(s4, Inches(1.55))

features = [
    ("Causal Graph Discovery", PURPLE,
     ["DoWhy + PC Algorithm", "Constraint-based causal discovery", "Fisher's Z-test for independence"]),
    ("Proxy Path Detection", RED,
     ["Backdoor criterion applied", "Direct & indirect paths found", "do-calculus effect estimates"]),
    ("Dual Protected Attributes", ACCENT,
     ["Age group (Young / Senior)", "Socioeconomic group (High / Low)", "Both audited simultaneously"]),
    ("Standard Fairness Benchmark", GREEN,
     ["AIF360 demographic parity", "Both DPDs computed & compared", "Shows model PASSING standard test"]),
    ("Impossibility Surface", RGBColor(0xF5,0x9E,0x0B),
     ["17-point threshold sweep", "Pareto frontier visualised", "Interactive slider in UI"]),
    ("Surgical Debiasing", RGBColor(0x06,0xB6,0xD4),
     ["Proxy variable removal", "Sample reweighting", "Before/after bias & accuracy"]),
    ("Gemini AI Reports", RGBColor(0xEC,0x48,0x99),
     ["Gemini 2.5 Flash powered", "Technical + plain-language", "Markdown formatted output"]),
    ("REST API", DARK_GRAY,
     ["FastAPI + Uvicorn", "POST /audit, POST /debias", "JSON responses, audit cache"]),
]

for i, (title, color, items) in enumerate(features):
    col = i % 4
    row = i // 4
    bx = Inches(0.55) + col * Inches(3.22)
    by = Inches(1.75) + row * Inches(2.55)
    box(s4, bx, by, Inches(3.05), Inches(2.35), fill=LIGHT_GRAY, line=color)
    # color top strip
    box(s4, bx, by, Inches(3.05), Inches(0.08), fill=color)
    txt(s4, title, bx + Inches(0.12), by + Inches(0.15),
        Inches(2.8), Inches(0.38), size=12, bold=True, color=color)
    for j, item in enumerate(items):
        txt(s4, "• " + item, bx + Inches(0.12),
            by + Inches(0.62) + j * Inches(0.48),
            Inches(2.8), Inches(0.42), size=11, color=DARK_GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — PROCESS FLOW
# ══════════════════════════════════════════════════════════════════════════════
s5 = slide()
accent_bar(s5)
section_tag(s5, "Slide 4", color=RGBColor(0xF5,0x9E,0x0B))
slide_number(s5, 4)
heading(s5, "Process Flow Diagram")
divider(s5, Inches(1.55))

steps = [
    ("1", "Upload\nDataset", "CSV or demo\ndataset loaded", PURPLE),
    ("2", "Standard\nAudit", "AIF360 DPD\nboth attributes", RGBColor(0x06,0xB6,0xD4)),
    ("3", "Causal\nGraph", "DoWhy builds\ncausal model", RED),
    ("4", "Proxy\nDetection", "Backdoor criterion\nfinds proxy paths", RGBColor(0xF5,0x9E,0x0B)),
    ("5", "Impossibility\nSurface", "Threshold sweep\nPareto frontier", GREEN),
    ("6", "Debiasing\n(optional)", "Proxy removal\n+ reweighting", ACCENT),
    ("7", "Gemini\nReport", "Plain-language\naudit report", RGBColor(0xEC,0x48,0x99)),
]

node_w = Inches(1.45)
node_h = Inches(1.55)
start_x = Inches(0.42)
node_y  = Inches(2.5)
gap     = Inches(0.28)

for i, (num, title, sub, color) in enumerate(steps):
    bx = start_x + i * (node_w + gap)
    # box
    box(s5, bx, node_y, node_w, node_h, fill=LIGHT_GRAY, line=color)
    box(s5, bx, node_y, node_w, Inches(0.07), fill=color)
    # number circle
    box(s5, bx + Inches(0.52), node_y + Inches(0.1),
        Inches(0.42), Inches(0.42), fill=color)
    txt(s5, num, bx + Inches(0.52), node_y + Inches(0.1),
        Inches(0.42), Inches(0.42), size=14, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
    txt(s5, title, bx + Inches(0.05), node_y + Inches(0.6),
        node_w - Inches(0.1), Inches(0.5),
        size=11, bold=True, color=color, align=PP_ALIGN.CENTER)
    txt(s5, sub, bx + Inches(0.05), node_y + Inches(1.1),
        node_w - Inches(0.1), Inches(0.5),
        size=9.5, color=GRAY, align=PP_ALIGN.CENTER)
    # arrow
    if i < len(steps) - 1:
        ax = bx + node_w + Inches(0.04)
        txt(s5, "→", ax, node_y + Inches(0.55),
            gap - Inches(0.04), Inches(0.45),
            size=16, bold=True, color=GRAY, align=PP_ALIGN.CENTER)

# Use-case note
box(s5, Inches(0.55), Inches(4.35), Inches(12.2), Inches(2.6),
    fill=RGBColor(0xF5,0xF3,0xFF), line=PURPLE)
txt(s5, "Key Use Cases", Inches(0.75), Inches(4.48),
    Inches(11.8), Inches(0.38), size=13, bold=True, color=PURPLE)
use_cases = [
    "HR Director uploads Q3 hiring dataset → Aition finds 3 proxy paths standard tools missed → Gemini report explains impact to board",
    "ML Engineer runs audit before model deployment → Impossibility slider shows fairness tradeoff → Surgical debiasing applied with 5.5% accuracy cost",
    "Compliance Officer generates PDPB audit trail → Both protected attributes documented → Debiasing strategy recorded with before/after metrics",
]
for j, uc in enumerate(use_cases):
    txt(s5, f"  {j+1}.  {uc}", Inches(0.75),
        Inches(4.92) + j * Inches(0.55),
        Inches(11.8), Inches(0.5), size=11, color=DARK_GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — WIREFRAMES / MOCK DIAGRAMS
# ══════════════════════════════════════════════════════════════════════════════
s6 = slide()
accent_bar(s6)
section_tag(s6, "Slide 5", color=RGBColor(0x06,0xB6,0xD4))
slide_number(s6, 5)
heading(s6, "Wireframes / Mock Diagrams")
divider(s6, Inches(1.55))

# Mock: Dashboard layout
dash_x, dash_y = Inches(0.55), Inches(1.75)
dash_w, dash_h = Inches(7.8), Inches(5.4)
box(s6, dash_x, dash_y, dash_w, dash_h,
    fill=RGBColor(0x05,0x08,0x18), line=PURPLE)
txt(s6, "DASHBOARD — Main View", dash_x + Inches(0.15),
    dash_y + Inches(0.1), Inches(7.5), Inches(0.3),
    size=9, bold=True, color=PURPLE_LT)

# Sidebar mock
box(s6, dash_x, dash_y, Inches(1.4), dash_h,
    fill=RGBColor(0x0A,0x0E,0x22), line=RGBColor(0x1A,0x1E,0x35))
for j, lbl in enumerate(["Overview","Fairness","Causal Graph","Debiasing","Reports"]):
    cy = dash_y + Inches(0.55) + j * Inches(0.52)
    col = PURPLE_LT if j == 0 else RGBColor(0x8A,0x8F,0xA8)
    txt(s6, ("▶ " if j==0 else "  ") + lbl,
        dash_x + Inches(0.08), cy, Inches(1.25), Inches(0.38),
        size=8, color=col, bold=(j==0))

# Alert banner mock
box(s6, dash_x + Inches(1.5), dash_y + Inches(0.45),
    Inches(6.1), Inches(0.55),
    fill=RGBColor(0x50,0x10,0x1E), line=RED)
txt(s6, "⚠  PROXY BIAS DETECTED  —  81 candidates affected",
    dash_x + Inches(1.65), dash_y + Inches(0.52),
    Inches(5.8), Inches(0.38), size=9, color=RGBColor(0xFF,0x7A,0x7A))

# Two metric cards
for j, (lbl, val, col) in enumerate([
    ("Standard Fairness (AIF360)", "FAIR ✓", GREEN),
    ("Causal Fairness (DoWhy)", "NOT FAIR ✗", RED),
]):
    cx = dash_x + Inches(1.5) + j * Inches(3.1)
    box(s6, cx, dash_y + Inches(1.1), Inches(2.9), Inches(1.3),
        fill=RGBColor(0x0E,0x13,0x30),
        line=col)
    txt(s6, lbl, cx + Inches(0.1), dash_y + Inches(1.18),
        Inches(2.7), Inches(0.3), size=8, color=RGBColor(0xC7,0xCA,0xE0))
    txt(s6, val, cx + Inches(0.1), dash_y + Inches(1.52),
        Inches(2.7), Inches(0.45), size=16, bold=True, color=col)

# Causal graph mock area
box(s6, dash_x + Inches(1.5), dash_y + Inches(2.55),
    Inches(6.1), Inches(2.65),
    fill=RGBColor(0x0E,0x13,0x30), line=RGBColor(0x1A,0x1E,0x35))
txt(s6, "Causal Graph", dash_x + Inches(1.65), dash_y + Inches(2.65),
    Inches(5.8), Inches(0.3), size=9, bold=True,
    color=RGBColor(0xE9,0xEB,0xF5))
# Simplified node sketch
for nx, ny, lbl, col in [
    (Inches(2.3), Inches(5.6), "Age\nGroup", PURPLE),
    (Inches(3.5), Inches(5.1), "Grad\nGap", RED),
    (Inches(3.5), Inches(6.0), "Emp\nGap", RED),
    (Inches(5.0), Inches(5.1), "Test\nScore", GREEN),
    (Inches(6.2), Inches(5.6), "Hired", GRAY),
]:
    box(s6, nx - Inches(0.28), ny - Inches(0.22),
        Inches(0.56), Inches(0.44), fill=col)
    txt(s6, lbl, nx - Inches(0.28), ny - Inches(0.22),
        Inches(0.56), Inches(0.44), size=7, color=WHITE,
        align=PP_ALIGN.CENTER)

# Right panel: annotations
ann_x = Inches(8.55)
txt(s6, "UI Annotations", ann_x, Inches(1.78),
    Inches(4.6), Inches(0.35), size=13, bold=True, color=DARK_GRAY)
annotations = [
    ("Sidebar", "Navigation with 5 sections"),
    ("Alert Banner", "Affected count + verdict"),
    ("Metric Cards", "Standard vs Causal verdict side-by-side"),
    ("Causal Graph", "SVG, colour-coded nodes & edges"),
    ("Impossibility Slider", "17-step threshold sweep"),
    ("Debiasing Panel", "Before/after bias & accuracy"),
    ("Report Section", "Gemini markdown rendered inline"),
]
for j, (k, v) in enumerate(annotations):
    by = Inches(2.2) + j * Inches(0.62)
    box(s6, ann_x, by, Inches(4.6), Inches(0.55),
        fill=LIGHT_GRAY, line=RGBColor(0xE5,0xE7,0xEB))
    txt(s6, k, ann_x + Inches(0.12), by + Inches(0.04),
        Inches(1.5), Inches(0.25), size=10, bold=True, color=PURPLE)
    txt(s6, v, ann_x + Inches(1.65), by + Inches(0.04),
        Inches(2.8), Inches(0.45), size=10, color=DARK_GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════════
s7 = slide()
accent_bar(s7)
section_tag(s7, "Slide 6", color=ACCENT)
slide_number(s7, 6)
heading(s7, "Architecture Diagram")
divider(s7, Inches(1.55))

layers = [
    ("FRONTEND  (React 18 + CRACO)", PURPLE,
     Inches(1.75), ["Dashboard.jsx", "CausalGraph.jsx", "Sidebar.jsx", "api.js"]),
    ("API GATEWAY  (FastAPI + Uvicorn)", RGBColor(0x06,0xB6,0xD4),
     Inches(3.1), ["POST /audit", "POST /audit/{id}/debias", "GET /health", "CORS middleware"]),
    ("AUDIT ENGINE  (Python)", RED,
     Inches(4.45), ["DoWhy CausalModel", "AIF360 BinaryLabelDataset", "Impossibility Surface", "Proxy Path Detection"]),
    ("DEBIASING + REPORT  (Python)", GREEN,
     Inches(5.8), ["sklearn LogisticRegression", "Sample Reweighting", "Gemini 2.5 Flash", "Fallback Report"]),
    ("DATA LAYER", DARK_GRAY,
     Inches(7.15), ["demo_hiring_dataset.csv", "In-memory audit cache", "generate_demo_dataset.py", ""]),
]

for title, color, y, items in layers:
    box(s7, Inches(0.55), y, Inches(12.2), Inches(1.15),
        fill=LIGHT_GRAY, line=color)
    box(s7, Inches(0.55), y, Inches(2.5), Inches(1.15), fill=color)
    txt(s7, title, Inches(0.65), y + Inches(0.38),
        Inches(2.3), Inches(0.45), size=10, bold=True,
        color=WHITE, align=PP_ALIGN.LEFT)
    for k, item in enumerate(items):
        if item:
            ix = Inches(3.25) + k * Inches(2.25)
            box(s7, ix, y + Inches(0.22), Inches(2.1), Inches(0.7),
                fill=WHITE, line=color)
            txt(s7, item, ix + Inches(0.08), y + Inches(0.3),
                Inches(1.95), Inches(0.5), size=10, color=DARK_GRAY)

# Arrows between layers
for ay in [Inches(2.9), Inches(4.25), Inches(5.6), Inches(6.95)]:
    txt(s7, "↕  HTTP / JSON  ↕", Inches(5.5), ay,
        Inches(2.3), Inches(0.3), size=9, color=GRAY,
        align=PP_ALIGN.CENTER, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — TECH STACK
# ══════════════════════════════════════════════════════════════════════════════
s8 = slide()
accent_bar(s8)
section_tag(s8, "Slide 7", color=GREEN)
slide_number(s8, 7)
heading(s8, "Technologies Used")
divider(s8, Inches(1.55))

tech = [
    ("Backend", PURPLE, [
        ("FastAPI + Uvicorn", "REST API framework, async, auto-docs"),
        ("DoWhy 0.11", "Causal graph discovery, PC algorithm, backdoor criterion"),
        ("AIF360 0.6", "Standard fairness metrics, BinaryLabelDatasetMetric"),
        ("scikit-learn", "LogisticRegression for debiasing baseline"),
        ("pandas + numpy", "Data manipulation and numerical computation"),
    ]),
    ("Frontend", RGBColor(0x06,0xB6,0xD4), [
        ("React 18", "Component-based UI, hooks, state management"),
        ("React Router v6", "Client-side routing"),
        ("Lucide React", "Icon library"),
        ("CRACO", "Create React App with custom webpack config"),
        ("Custom SVG", "CausalGraph rendered without D3 dependency"),
    ]),
    ("AI & Cloud", RED, [
        ("Google Gemini 2.5 Flash", "Plain-language audit report generation"),
        ("google-genai SDK", "Python client for Gemini API"),
        ("Vercel", "Frontend deployment (static build)"),
        ("python-dotenv", "Environment variable management"),
        ("Hypothesis", "Property-based testing for backend"),
    ]),
]

for i, (cat, color, items) in enumerate(tech):
    bx = Inches(0.55) + i * Inches(4.25)
    box(s8, bx, Inches(1.75), Inches(4.05), Inches(5.4),
        fill=LIGHT_GRAY, line=color)
    box(s8, bx, Inches(1.75), Inches(4.05), Inches(0.07), fill=color)
    txt(s8, cat, bx + Inches(0.15), Inches(1.85),
        Inches(3.75), Inches(0.38), size=14, bold=True, color=color)
    for j, (tech_name, desc) in enumerate(items):
        ty = Inches(2.35) + j * Inches(0.88)
        box(s8, bx + Inches(0.12), ty, Inches(3.8), Inches(0.78),
            fill=WHITE, line=RGBColor(0xE5,0xE7,0xEB))
        txt(s8, tech_name, bx + Inches(0.25), ty + Inches(0.06),
            Inches(3.55), Inches(0.3), size=11, bold=True, color=DARK_GRAY)
        txt(s8, desc, bx + Inches(0.25), ty + Inches(0.38),
            Inches(3.55), Inches(0.35), size=9.5, color=GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — ESTIMATED IMPLEMENTATION COST
# ══════════════════════════════════════════════════════════════════════════════
s9 = slide()
accent_bar(s9)
section_tag(s9, "Slide 8", color=RGBColor(0xF5,0x9E,0x0B))
slide_number(s9, 8)
heading(s9, "Estimated Implementation Cost")
divider(s9, Inches(1.55))

cost_items = [
    ("Google Gemini API", "~$0.075 / 1M input tokens\n~$0.30 / 1M output tokens",
     "~$5–15/month at demo scale\n~$50–200/month at production scale",
     RGBColor(0xF5,0x9E,0x0B)),
    ("Vercel (Frontend)", "Free tier: unlimited personal projects\nPro: $20/month/member",
     "Free for prototype\n$20/month for team deployment",
     PURPLE),
    ("Cloud Run (Backend)", "GCP: $0.00002400/vCPU-second\n$0.00000250/GB-second",
     "~$0 at demo scale (free tier)\n~$30–80/month at production",
     RGBColor(0x06,0xB6,0xD4)),
    ("Development (MVP)", "1 developer, 3 months\nPart-time student project",
     "₹0 (student project)\n~$15,000 commercial equivalent",
     GREEN),
]

for i, (item, detail, cost, color) in enumerate(cost_items):
    col = i % 2
    row = i // 2
    bx = Inches(0.55) + col * Inches(6.3)
    by = Inches(1.85) + row * Inches(2.55)
    box(s9, bx, by, Inches(6.0), Inches(2.3), fill=LIGHT_GRAY, line=color)
    box(s9, bx, by, Inches(6.0), Inches(0.07), fill=color)
    txt(s9, item, bx + Inches(0.18), by + Inches(0.15),
        Inches(5.65), Inches(0.38), size=14, bold=True, color=color)
    txt(s9, detail, bx + Inches(0.18), by + Inches(0.62),
        Inches(5.65), Inches(0.7), size=11, color=DARK_GRAY)
    box(s9, bx + Inches(0.18), by + Inches(1.38),
        Inches(5.65), Inches(0.72), fill=WHITE, line=color)
    txt(s9, "💰  " + cost, bx + Inches(0.3), by + Inches(1.45),
        Inches(5.4), Inches(0.6), size=11, bold=True, color=color)

txt(s9, "Total estimated monthly cost at production scale: ~$100–300/month",
    Inches(0.55), Inches(7.05), Inches(12.2), Inches(0.35),
    size=12, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — MVP SNAPSHOTS
# ══════════════════════════════════════════════════════════════════════════════
s10 = slide()
accent_bar(s10)
section_tag(s10, "Slide 9", color=RED)
slide_number(s10, 9)
heading(s10, "Snapshots of the MVP")
divider(s10, Inches(1.55))

# Left: audit result mock
mock_x, mock_y = Inches(0.55), Inches(1.75)
mock_w, mock_h = Inches(8.0), Inches(5.5)
box(s10, mock_x, mock_y, mock_w, mock_h,
    fill=RGBColor(0x05,0x08,0x18), line=PURPLE)

# Alert
box(s10, mock_x + Inches(0.15), mock_y + Inches(0.2),
    mock_w - Inches(0.3), Inches(0.55),
    fill=RGBColor(0x50,0x10,0x1E), line=RED)
txt(s10, "⚠  PROXY DISCRIMINATION DETECTED  —  81 candidates affected",
    mock_x + Inches(0.3), mock_y + Inches(0.28),
    mock_w - Inches(0.5), Inches(0.35),
    size=9, color=RGBColor(0xFF,0x7A,0x7A))

# Metric cards
for j, (lbl, val, col) in enumerate([
    ("Standard Fairness", "FAIR ✓  DPD: 0.0809", GREEN),
    ("Causal Fairness", "NOT FAIR ✗  3 paths", RED),
]):
    cx = mock_x + Inches(0.15) + j * Inches(3.95)
    box(s10, cx, mock_y + Inches(0.9), Inches(3.75), Inches(1.1),
        fill=RGBColor(0x0E,0x13,0x30), line=col)
    txt(s10, lbl, cx + Inches(0.12), mock_y + Inches(0.98),
        Inches(3.5), Inches(0.28), size=9, color=RGBColor(0xC7,0xCA,0xE0))
    txt(s10, val, cx + Inches(0.12), mock_y + Inches(1.3),
        Inches(3.5), Inches(0.45), size=14, bold=True, color=col)

# Proxy paths
box(s10, mock_x + Inches(0.15), mock_y + Inches(2.15),
    mock_w - Inches(0.3), Inches(1.5),
    fill=RGBColor(0x0E,0x13,0x30), line=PURPLE)
txt(s10, "Proxy Paths Detected by Aition:",
    mock_x + Inches(0.3), mock_y + Inches(2.22),
    mock_w - Inches(0.5), Inches(0.3),
    size=9, bold=True, color=PURPLE_LT)
for j, path in enumerate([
    "age_group  →  college_graduation_year_gap  →  hired   (effect: 0.298)",
    "age_group  →  employment_gap  →  hired   (effect: 0.152)",
    "socioeconomic_group  →  neighborhood_quality  →  hired   (effect: 0.401)",
]):
    txt(s10, "🔴  " + path,
        mock_x + Inches(0.3), mock_y + Inches(2.58) + j * Inches(0.33),
        mock_w - Inches(0.5), Inches(0.3),
        size=9, color=RGBColor(0xFF,0x9A,0x9A))

# Debiasing result
box(s10, mock_x + Inches(0.15), mock_y + Inches(3.82),
    mock_w - Inches(0.3), Inches(1.45),
    fill=RGBColor(0x0E,0x13,0x30), line=GREEN)
txt(s10, "Surgical Debiasing Result:",
    mock_x + Inches(0.3), mock_y + Inches(3.9),
    mock_w - Inches(0.5), Inches(0.3),
    size=9, bold=True, color=RGBColor(0x34,0xD3,0x99))
for j, (k, v) in enumerate([
    ("Bias Index", "0.0603  →  0.0269   (65.4% reduction)"),
    ("Accuracy",   "93.8%  →  88.3%   (5.5% cost)"),
    ("Removed",    "college_graduation_year_gap, neighborhood_quality"),
]):
    txt(s10, f"{k}:  {v}",
        mock_x + Inches(0.3), mock_y + Inches(4.28) + j * Inches(0.3),
        mock_w - Inches(0.5), Inches(0.28),
        size=9, color=RGBColor(0x6E,0xE7,0xB7))

# Right: key numbers
rp_x = Inches(8.75)
txt(s10, "Key MVP Numbers", rp_x, Inches(1.78),
    Inches(4.3), Inches(0.35), size=13, bold=True, color=DARK_GRAY)

for j, (val, lbl, col) in enumerate([
    ("3 / 0", "Proxy paths: Aition vs Standard", PURPLE),
    ("81", "Candidates wrongly affected", RED),
    ("65.4%", "Bias reduction after debiasing", GREEN),
    ("5.5%", "Accuracy cost of debiasing", RGBColor(0xF5,0x9E,0x0B)),
    ("2", "Protected attributes audited", RGBColor(0x06,0xB6,0xD4)),
    ("0.0997", "DPD — passes standard threshold", GRAY),
]):
    by = Inches(2.2) + j * Inches(0.82)
    box(s10, rp_x, by, Inches(4.3), Inches(0.72),
        fill=LIGHT_GRAY, line=col)
    txt(s10, val, rp_x + Inches(0.15), by + Inches(0.06),
        Inches(1.4), Inches(0.58), size=18, bold=True, color=col)
    txt(s10, lbl, rp_x + Inches(1.6), by + Inches(0.18),
        Inches(2.55), Inches(0.38), size=10, color=DARK_GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — FUTURE DEVELOPMENT
# ══════════════════════════════════════════════════════════════════════════════
s11 = slide()
accent_bar(s11)
section_tag(s11, "Slide 10", color=ACCENT)
slide_number(s11, 10)
heading(s11, "Future Development")
divider(s11, Inches(1.55))

phases = [
    ("Phase 1\n(3 months)", PURPLE, [
        "Unstructured input auditing (CV text, voice screening)",
        "Multi-language Gemini reports (Hindi, Tamil, Telugu)",
        "Firebase audit trail — immutable, versioned",
        "User authentication + organisation management",
    ]),
    ("Phase 2\n(6 months)", RGBColor(0x06,0xB6,0xD4), [
        "Real-time bias drift monitoring across model versions",
        "Support for custom protected attributes (any column)",
        "Integration with ATS platforms (Workday, Greenhouse)",
        "PDPB 2023 compliance certificate generation",
    ]),
    ("Phase 3\n(12 months)", GREEN, [
        "Open-source audit standard for Indian AI systems",
        "National fairness benchmark dataset (India-specific)",
        "Credit scoring and healthcare triage audit modules",
        "Flutter mobile app for field compliance officers",
    ]),
]

for i, (phase, color, items) in enumerate(phases):
    bx = Inches(0.55) + i * Inches(4.25)
    box(s11, bx, Inches(1.75), Inches(4.05), Inches(4.9),
        fill=LIGHT_GRAY, line=color)
    box(s11, bx, Inches(1.75), Inches(4.05), Inches(0.07), fill=color)
    txt(s11, phase, bx + Inches(0.15), Inches(1.85),
        Inches(3.75), Inches(0.55), size=14, bold=True, color=color)
    for j, item in enumerate(items):
        ty = Inches(2.55) + j * Inches(0.72)
        box(s11, bx + Inches(0.12), ty, Inches(3.8), Inches(0.62),
            fill=WHITE, line=RGBColor(0xE5,0xE7,0xEB))
        txt(s11, "→  " + item, bx + Inches(0.22), ty + Inches(0.08),
            Inches(3.6), Inches(0.5), size=11, color=DARK_GRAY)

# SDG alignment
box(s11, Inches(0.55), Inches(6.85), Inches(12.2), Inches(0.5),
    fill=RGBColor(0xF5,0xF3,0xFF), line=PURPLE)
txt(s11,
    "SDG Alignment:  SDG 10 Reduced Inequalities  ·  SDG 8 Decent Work  ·  SDG 16 Justice & Institutions  ·  SDG 5 Gender Equality",
    Inches(0.75), Inches(6.9), Inches(11.8), Inches(0.38),
    size=11, color=PURPLE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════════════════
prs.save("Aition_Prototype_Deck.pptx")
print("✅  Saved: Aition_Prototype_Deck.pptx")
